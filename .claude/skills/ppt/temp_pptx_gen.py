#!/usr/bin/env python3
# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os
from datetime import datetime

ONEDRIVE_PATH = r"C:\Users\민준\OneDrive - 주식회사 도리당\data\report\PowerPoint"
os.makedirs(ONEDRIVE_PATH, exist_ok=True)
FILE_NAME = "AI_데이터기반_의사결정_2_애니메이션없는_ppt.pptx"
OUTPUT_PATH = os.path.join(ONEDRIVE_PATH, FILE_NAME)

DARK_BG   = RGBColor(13, 13, 13)
MINT      = RGBColor(0, 229, 160)
DARK_MINT = RGBColor(0, 160, 110)
WHITE     = RGBColor(255, 255, 255)
LGRAY     = RGBColor(180, 180, 180)
CARD_BG   = RGBColor(30, 30, 30)
CARD_BG2  = RGBColor(20, 45, 35)
RED       = RGBColor(220, 80, 80)
RED_DARK  = RGBColor(180, 50, 50)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


# ── 공통 유틸 ──────────────────────────────────────────
def bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG

def txbox(slide, l, t, w, h, text, size,
          bold=False, color=None, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color if color else WHITE
    return tb

def rect(slide, l, t, w, h, fill_rgb, line_rgb=None, line_width_pt=1):
    sh = slide.shapes.add_shape(1,
        Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill_rgb
    if line_rgb:
        sh.line.color.rgb = line_rgb
        sh.line.width = Pt(line_width_pt)
    else:
        sh.line.fill.background()
    return sh

def title_bar(slide, title_text, font_size=34):
    rect(slide, 0, 0, 13.333, 0.85, MINT)
    txbox(slide, 0.4, 0.09, 12.5, 0.66, title_text, font_size,
          bold=True, color=DARK_BG)

def set_notes(slide, script):
    slide.notes_slide.notes_text_frame.text = script


# ── 슬라이드 함수들 ─────────────────────────────────────

def slide_cover(title, subtitle, date_str):
    s = prs.slides.add_slide(BLANK)
    bg(s)
    rect(s, 0, 0, 13.333, 1.8, MINT)
    rect(s, 12.5, 0, 0.833, 1.8, DARK_MINT)
    rect(s, 0, 7.2, 13.333, 0.12, MINT)
    txbox(s, 0.5, 2.1, 12, 1.6, title, 52,
          bold=True, align=PP_ALIGN.CENTER)
    txbox(s, 0.5, 3.9, 12, 0.8, subtitle, 22,
          color=LGRAY, align=PP_ALIGN.CENTER)
    txbox(s, 0.5, 6.5, 12, 0.6, date_str, 15,
          color=LGRAY, align=PP_ALIGN.CENTER)


def slide_content(title, content_text, script):
    s = prs.slides.add_slide(BLANK)
    bg(s)
    title_bar(s, title)
    txbox(s, 0.8, 1.1, 11.7, 5.9, content_text, 19, wrap=True)
    set_notes(s, script)


def slide_two_col(title, l_title, l_items, r_title, r_items, script):
    s = prs.slides.add_slide(BLANK)
    bg(s)
    title_bar(s, title)
    rect(s, 0.4, 1.05, 5.9, 5.9, CARD_BG)
    txbox(s, 0.7, 1.2, 5.3, 0.55, l_title, 19, bold=True, color=LGRAY)
    txbox(s, 0.7, 1.85, 5.3, 4.8, l_items, 17, wrap=True)
    rect(s, 7.0, 1.05, 5.9, 5.9, CARD_BG2, line_rgb=MINT, line_width_pt=2)
    txbox(s, 7.3, 1.2, 5.3, 0.55, r_title, 19, bold=True, color=MINT)
    txbox(s, 7.3, 1.85, 5.3, 4.8, r_items, 17, wrap=True)
    set_notes(s, script)


def slide_flow(title, steps, script):
    s = prs.slides.add_slide(BLANK)
    bg(s)
    title_bar(s, title)
    n = len(steps)
    box_w  = 1.6
    gap    = 0.22
    arr_w  = 0.28
    total  = n * box_w + (n - 1) * (gap + arr_w)
    sx     = (13.333 - total) / 2
    top    = 2.1
    bh     = 3.8

    for i, (num, label, sub) in enumerate(steps):
        x = sx + i * (box_w + gap + arr_w)
        rect(s, x, top, box_w, bh, CARD_BG, line_rgb=MINT, line_width_pt=1.5)
        txbox(s, x + 0.05, top + 0.15, box_w - 0.1, 0.55,
              num, 24, bold=True, color=MINT, align=PP_ALIGN.CENTER)
        txbox(s, x + 0.05, top + 0.82, box_w - 0.1, 1.2,
              label, 14, bold=True, align=PP_ALIGN.CENTER)
        if sub:
            txbox(s, x + 0.05, top + 2.1, box_w - 0.1, 1.5,
                  sub, 11, color=LGRAY, align=PP_ALIGN.CENTER)
        if i < n - 1:
            ax = x + box_w + gap / 2
            txbox(s, ax, top + 1.6, arr_w + 0.1, 0.55,
                  "→", 22, color=MINT, align=PP_ALIGN.CENTER)
    set_notes(s, script)


def slide_case(title, problem, goal, solution, script):
    s = prs.slides.add_slide(BLANK)
    bg(s)
    title_bar(s, title)
    rect(s, 0.35, 1.05, 5.9, 2.85, CARD_BG)
    txbox(s, 0.55, 1.12, 0.9, 0.55, "①", 22, bold=True, color=MINT)
    txbox(s, 1.3,  1.12, 4.7, 0.55, "문제 정의", 17, bold=True, color=LGRAY)
    txbox(s, 0.55, 1.72, 5.5, 2.0, problem, 15, wrap=True)

    rect(s, 6.95, 1.05, 5.9, 2.85, CARD_BG)
    txbox(s, 7.15, 1.12, 0.9, 0.55, "②", 22, bold=True, color=MINT)
    txbox(s, 7.9,  1.12, 4.7, 0.55, "목표 설정", 17, bold=True, color=LGRAY)
    txbox(s, 7.15, 1.72, 5.5, 2.0, goal, 15, wrap=True)

    rect(s, 0.35, 4.1, 12.5, 2.9, CARD_BG2, line_rgb=MINT, line_width_pt=1.5)
    txbox(s, 0.55, 4.17, 0.9, 0.55, "③", 22, bold=True, color=MINT)
    txbox(s, 1.3,  4.17, 11.0, 0.55, "해결 방법", 17, bold=True, color=MINT)
    txbox(s, 0.55, 4.77, 12.0, 2.1, solution, 15, wrap=True)
    set_notes(s, script)


def slide_kpi(title, cards, script):
    s = prs.slides.add_slide(BLANK)
    bg(s)
    title_bar(s, title)
    cw = 3.8
    ch = 5.3
    xs = [0.35, 4.73, 9.12]
    for (icon, heading, body), x in zip(cards, xs):
        rect(s, x, 1.1, cw, ch, CARD_BG, line_rgb=MINT, line_width_pt=1.5)
        txbox(s, x + 0.1, 1.22, cw - 0.2, 0.9,
              icon, 34, align=PP_ALIGN.CENTER)
        txbox(s, x + 0.1, 2.18, cw - 0.2, 0.72,
              heading, 16, bold=True, color=MINT, align=PP_ALIGN.CENTER)
        txbox(s, x + 0.15, 2.96, cw - 0.3, 3.3,
              body, 14, color=LGRAY, wrap=True, align=PP_ALIGN.CENTER)
    set_notes(s, script)


def slide_quote(title, quote, points, script):
    s = prs.slides.add_slide(BLANK)
    bg(s)
    title_bar(s, title)
    rect(s, 0.5, 1.05, 12.3, 1.55, CARD_BG2, line_rgb=MINT, line_width_pt=2)
    txbox(s, 0.8, 1.12, 11.8, 1.38,
          f'"{quote}"', 22, bold=True, color=MINT, align=PP_ALIGN.CENTER)
    cw = 3.7
    xs = [0.5, 4.77, 9.05]
    for (heading, body), x in zip(points, xs):
        rect(s, x, 2.8, cw, 4.25, CARD_BG)
        txbox(s, x + 0.15, 2.9, cw - 0.3, 0.62, heading, 15, bold=True, color=MINT)
        txbox(s, x + 0.15, 3.57, cw - 0.3, 3.3, body, 13, color=LGRAY, wrap=True)
    set_notes(s, script)


def slide_conclusion(title, main_msg, bad_text, good_text, script):
    s = prs.slides.add_slide(BLANK)
    bg(s)
    title_bar(s, title)
    rect(s, 0.5, 1.05, 12.3, 1.85, CARD_BG2, line_rgb=MINT, line_width_pt=2)
    txbox(s, 0.65, 1.15, 12.1, 1.65,
          main_msg, 26, bold=True, color=MINT, align=PP_ALIGN.CENTER)
    rect(s, 0.5, 3.15, 5.9, 3.75, CARD_BG, line_rgb=RED_DARK, line_width_pt=1.5)
    txbox(s, 0.7, 3.25, 5.5, 0.58, "✕  많이 사용하는 것", 17, bold=True, color=RED)
    txbox(s, 0.7, 3.88, 5.5, 2.9, bad_text, 14, color=LGRAY, wrap=True)
    rect(s, 7.0, 3.15, 5.9, 3.75, CARD_BG2, line_rgb=MINT, line_width_pt=2)
    txbox(s, 7.2, 3.25, 5.5, 0.58, "✓  실행까지 연결하는 것", 17, bold=True, color=MINT)
    txbox(s, 7.2, 3.88, 5.5, 2.9, good_text, 14, wrap=True)
    set_notes(s, script)


# ═══════════════════════════════════════════════════════
# 슬라이드 생성
# ═══════════════════════════════════════════════════════
today = datetime.now().strftime("%Y년 %m월 %d일")

# 1. 표지
slide_cover("AI 데이터 기반 의사결정",
            "반복 업무를 빠르게 처리하는 방법  |  직원 교육", today)

# 2. 교육 목적
slide_content("교육 목적",
    "이번 교육은 AI 자체를 배우는 교육이 아닙니다.\n\n"
    "• 보고서 정리 시간 단축\n"
    "• 데이터 분석 속도 향상\n"
    "• 인수인계 시간 감소\n\n"
    "처럼 내 업무를 더 빠르고 쉽게 처리하는 방법을\n"
    "실제 사례 중심으로 경험하는 교육입니다.",
    "이번 교육은 AI 자체를 배우는 교육이 아닙니다.\n"
    "보고서 정리, 데이터 분석, 인수인계 같은 반복 업무를\n"
    "조금 더 빠르고 쉽게 처리하는 방법을\n"
    "실제 사례 중심으로 같이 보는 시간입니다.")

# 3. DX vs AX
slide_two_col(
    "DX vs AX — 무엇이 다른가?",
    "DX (Digital Transformation)",
    "데이터를 빠르게 보고 공유하는 구조\n\n"
    "• 데이터 수집\n"
    "• 대시보드\n"
    "• 알림 시스템\n\n"
    "\"데이터를 보는 속도\"를 높이는 것",
    "AX (AI Transformation)",
    "AI를 활용해 의사결정과 실행까지 연결\n\n"
    "• AI가 데이터 분석\n"
    "• AI가 문제 원인 해석\n"
    "• AI가 액션 제안\n\n"
    "\"의사결정 속도\"를 높이는 것",
    "DX는 데이터를 빠르게 보고 공유하는 구조입니다.\n"
    "예를 들면 데이터 수집, 대시보드, 알림 시스템 같은 것들입니다.\n"
    "AX는 여기서 한 단계 더 나아가\n"
    "AI가 데이터를 분석하고, 문제 원인을 해석하고,\n"
    "액션까지 제안하는 구조입니다.\n"
    "결국 핵심은 AI를 활용해서 실행까지 연결하는 것입니다.")

# 4. 활용 방법 (프레임워크)
slide_flow(
    "AI 활용 기본 프레임워크",
    [
        ("①", "문제 정의", "해결해야 할\n핵심 문제"),
        ("②", "목표 설정", "원하는\n결과 명확화"),
        ("③", "데이터 첨부", "프로젝트 연결\n+ 데이터"),
        ("④", "AI 질문", "명확하게\n요청"),
        ("⑤", "의사결정", "AI 결과를\n판단"),
        ("⑥", "실행", "바로\n적용"),
    ],
    "오늘은 복잡하게 하지 않고 이 흐름 하나만 기억하시면 됩니다.\n"
    "문제 정의 → 목표 설정 → 프로젝트 연결 + 데이터 첨부\n"
    "→ AI 질문 → 의사결정 → 실행\n"
    "중요한 건 AI에게 그냥 질문하는 게 아니라\n"
    "업무 흐름 안에서 활용하는 것입니다.")

# 5. 사례1: 마케팅팀
slide_case(
    "사례1: 마케팅팀 — 담당자 변경 시 공백 최소화",
    "담당자가 변경되면 이전 광고 운영 내용과\n"
    "매장 이슈 파악이 어려움\n\n→ 정보 공백 발생",
    "담당자 변경 시에도\n"
    "광고 운영 흐름과 핵심 이슈를\n빠르게 파악",
    "[마케팅 요청사항] 프로젝트 연결\n"
    "+ 최근 광고 운영 내역 & 매장별 특이사항 첨부\n\n"
    "→ \"광고 운영 내역과 매장별 이슈를 정리해줘\"\n"
    "→ 인수인계 시간을 크게 줄일 수 있습니다",
    "예를 들면 마케팅팀은 담당자가 바뀌면\n"
    "이전 광고 운영 내용이나 매장 이슈 파악이 어려운 경우가 많습니다.\n"
    "이럴 때 프로젝트와 데이터를 연결한 뒤\n"
    "AI에게 광고 운영 내역이나 특이사항을 정리하게 하면\n"
    "인수인계 시간을 크게 줄일 수 있습니다.")

# 6. 사례2: 영업관리부
slide_case(
    "사례2: 영업관리부 — 점주 대응 전략 수립",
    "송파점 점주가 매출 관련\n"
    "부정적인 의견을 반복적으로 제기\n\n→ 대응 전략 필요",
    "점주 설득 및\n"
    "실행 가능한 개선 방향 도출",
    "[서울_송파점] 프로젝트 연결 + 최근 매출 데이터 첨부\n\n"
    "→ \"최근 방문일지 기준 이슈 정리, 점주 성향/반복 불만 분석,\n"
    "   매출 대응 방향 제안해줘\"\n"
    "→ 단순 요약이 아닌 실제 대응 방향까지 정리",
    "영업관리부는 방문일지 활용 사례입니다.\n"
    "예를 들어 송파점 점주가 매출 관련 부정적인 의견을 반복적으로 이야기하고 있다고 하면\n"
    "최근 매출 데이터와 방문일지를 같이 넣고\n"
    "'점주의 성향과 반복 불만을 분석해서 대응 방향을 제안해줘'\n"
    "라고 질문할 수 있습니다.\n"
    "그러면 단순 요약이 아니라 실제 대응 방향까지 같이 정리할 수 있습니다.")

# 7. 프레임워크 재강조
slide_flow(
    "핵심 프레임워크 — 모든 업무에 동일하게 적용",
    [
        ("①", "문제 정의", ""),
        ("②", "목표 설정", ""),
        ("③", "데이터 첨부", ""),
        ("④", "AI 질문", ""),
        ("⑤", "의사결정", ""),
        ("⑥", "실행", ""),
    ],
    "결국 핵심 구조는 동일합니다.\n"
    "문제 정의 → 목표 설정 → 프로젝트 연결 + 데이터 첨부\n"
    "→ AI 질문 → 의사결정 → 실행\n"
    "이 흐름 안에서 AI를 활용하는 게 중요합니다.")

# 8. 실전팁
slide_quote(
    "실전팁",
    "일단 업무에 붙여서 자주 써보세요",
    [
        ("어떤 방식으로 질문?",
         "어떤 방식으로 질문해야\n"
         "원하는 답이 나오는지\n"
         "반복하면서 감각이 생깁니다"),
        ("어떤 정보를 넣어야?",
         "어떤 데이터를 함께\n"
         "제공해야 좋은 결과가\n"
         "나오는지 알게 됩니다"),
        ("어떤 형태로 요청?",
         "어떤 형태로 요청해야\n"
         "바로 쓸 수 있는 결과가\n"
         "나오는지 체감합니다"),
    ],
    "좋은 프롬프트를 외우는 것보다\n"
    "일단 업무에 많이 붙여서 써보는 걸 추천드립니다.\n"
    "계속 사용하다 보면\n"
    "어떤 정보를 넣어야 하는지,\n"
    "어떻게 질문해야 원하는 결과가 나오는지\n"
    "자연스럽게 감각이 생기기 시작합니다.")

# 9. 기대효과
slide_kpi(
    "기대효과",
    [
        ("⏱", "반복 업무 시간 단축",
         "보고서 정리, 데이터 분석,\n인수인계 등 반복 업무를\nAI로 빠르게 처리"),
        ("📊", "의사결정 속도 향상",
         "데이터 기반으로\n더 빠르고 정확한\n의사결정 가능"),
        ("🚀", "실행 문화 강화",
         "부서별 활용 사례 확산\n및 AI 활용\n실행 문화 정착"),
    ],
    "반복 업무 정리 시간은 줄이고\n"
    "데이터 기반 의사결정 속도는 높이고\n"
    "부서별 활용 사례도 점점 확산될 수 있을 것으로 기대하고 있습니다.")

# 10. 결론
slide_conclusion(
    "결론",
    "AI는 의사결정을 빠르게 만드는 도구입니다",
    "사용 횟수만 늘리는 것은\n진정한 활용이 아닙니다.\n\n"
    "단순히 AI를 켜두는 것으로는\n업무가 개선되지 않습니다.",
    "문제 정의 → 데이터 첨부\n→ AI 질문 → 의사결정 → 실행\n\n"
    "이 흐름을 실제 업무에 연결할 때\n진정한 효과가 생깁니다.",
    "AI는 단순히 데이터를 정리하는 도구가 아니라\n"
    "의사결정을 빠르게 만드는 도구라고 생각합니다.\n"
    "중요한 건 많이 사용하는 것보다\n"
    "실제 업무에 활용해서 실행까지 연결하는 것입니다.")

# 저장
prs.save(OUTPUT_PATH)
print(f"OK | {OUTPUT_PATH}")
print(f"Slides: {len(prs.slides)}")
